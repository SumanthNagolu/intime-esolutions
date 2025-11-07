#!/usr/bin/env python3
"""
Extract text content from PPTs and PDFs in data/ folder
Save to JSON for ingestion into database with embeddings
"""

import os
import json
from pathlib import Path
from pptx import Presentation
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    print("⚠️  PyPDF2 not installed. PDF extraction will be skipped.")
    print("   Install with: pip install PyPDF2")
    PDF_AVAILABLE = False

def extract_ppt_slides(ppt_path):
    """Extract text from each slide"""
    try:
        prs = Presentation(ppt_path)
        slides = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                content = '\n'.join(slide_text)
                slides.append({
                    'number': idx,
                    'content': content,
                    'word_count': len(content.split())
                })
        
        return slides
    except Exception as e:
        print(f"    ⚠️  Error extracting {ppt_path.name}: {e}")
        return []

def extract_pdf_pages(pdf_path):
    """Extract text from PDF pages"""
    if not PDF_AVAILABLE:
        return []
    
    try:
        pages = []
        
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            
            for idx, page in enumerate(reader.pages, start=1):
                text = page.extract_text().strip()
                if text:
                    pages.append({
                        'number': idx,
                        'content': text,
                        'word_count': len(text.split())
                    })
        
        return pages
    except Exception as e:
        print(f"    ⚠️  Error extracting {pdf_path.name}: {e}")
        return []

def process_topic_folder(topic_path):
    """Process all files in a topic folder"""
    topic_code = topic_path.name
    content = {
        'topic_code': topic_code,
        'slides': [],
        'assignment_pages': [],
        'demo_files': []
    }
    
    for file in topic_path.iterdir():
        if file.suffix.lower() in ['.pptx', '.ppt']:
            print(f"    📄 Extracting {file.name}...")
            slides = extract_ppt_slides(file)
            content['slides'].extend(slides)
        
        elif file.suffix.lower() == '.pdf':
            if 'assignment' in file.name.lower() or 'solution' in file.name.lower():
                print(f"    📄 Extracting {file.name}...")
                pages = extract_pdf_pages(file)
                content['assignment_pages'].extend(pages)
        
        elif file.suffix.lower() in ['.mp4', '.mov', '.avi']:
            # Store demo file reference
            content['demo_files'].append({
                'filename': file.name,
                'path': str(file.relative_to(Path('data'))),
                'size_mb': round(file.stat().st_size / (1024 * 1024), 2)
            })
    
    return content

def main():
    data_dir = Path('data')
    output_dir = Path('content/extracted')
    output_dir.mkdir(exist_ok=True, parents=True)
    
    if not data_dir.exists():
        print(f"❌ Data directory not found: {data_dir}")
        return
    
    print("🚀 Starting content extraction...\n")
    
    total_slides = 0
    total_pages = 0
    total_topics = 0
    
    # Process each product folder
    for product_folder in sorted(data_dir.iterdir()):
        if not product_folder.is_dir():
            continue
        
        print(f"📁 Processing {product_folder.name}/")
        
        for topic_folder in sorted(product_folder.iterdir()):
            if not topic_folder.is_dir():
                continue
            
            print(f"  📚 {topic_folder.name}")
            content = process_topic_folder(topic_folder)
            
            if content['slides'] or content['assignment_pages']:
                # Save to JSON
                output_file = output_dir / f"{topic_folder.name}.json"
                with open(output_file, 'w', encoding='utf-8') as f:
                    json.dump(content, f, indent=2, ensure_ascii=False)
                
                total_slides += len(content['slides'])
                total_pages += len(content['assignment_pages'])
                total_topics += 1
                
                print(f"    ✅ {len(content['slides'])} slides, "
                      f"{len(content['assignment_pages'])} pages, "
                      f"{len(content['demo_files'])} demos")
            else:
                print(f"    ⚠️  No content extracted")
    
    print(f"\n{'='*60}")
    print(f"✅ Content extraction complete!")
    print(f"{'='*60}")
    print(f"📊 Statistics:")
    print(f"   • Topics processed: {total_topics}")
    print(f"   • Slides extracted: {total_slides}")
    print(f"   • Assignment pages: {total_pages}")
    print(f"   • Total chunks: {total_slides + total_pages}")
    print(f"\n📁 Output directory: {output_dir}")
    print(f"\n🎯 Next step: Run 'npx tsx scripts/generate-embeddings.ts'")

if __name__ == '__main__':
    main()

